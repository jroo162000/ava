#!/usr/bin/env python3
"""AVA rich document generator. Usage: ava_filegen.py <path> <format>
Content is read from stdin (UTF-8). Emits a JSON status line on stdout.
Supports: docx (python-docx), xlsx (openpyxl), pptx (python-pptx)."""
import sys, os, json

def main():
    try:
        sys.stdin.reconfigure(encoding="utf-8")
    except Exception:
        pass
    if len(sys.argv) < 3:
        print(json.dumps({"ok": False, "error": "usage: ava_filegen.py <path> <format>"})); return
    path = sys.argv[1]
    fmt = sys.argv[2].lower().lstrip(".")
    content = sys.stdin.read()
    try:
        d = os.path.dirname(path)
        if d:
            os.makedirs(d, exist_ok=True)

        if fmt == "docx":
            import docx
            doc = docx.Document()
            lines = content.split("\n") if content else [""]
            # First non-empty line becomes a heading if it looks like a title
            for i, line in enumerate(lines):
                if i == 0 and line.strip():
                    doc.add_heading(line.strip(), level=1)
                else:
                    doc.add_paragraph(line)
            doc.save(path)

        elif fmt == "xlsx":
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            for r, line in enumerate(content.split("\n"), start=1):
                cells = line.split("\t") if "\t" in line else line.split(",")
                for c, val in enumerate(cells, start=1):
                    ws.cell(row=r, column=c, value=val)
            wb.save(path)

        elif fmt == "pptx":
            import pptx
            prs = pptx.Presentation()
            blocks = content.split("\n\n") if "\n\n" in content else [content]
            for b in blocks:
                lines = [x for x in b.split("\n")]
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                try:
                    slide.shapes.title.text = lines[0] if lines and lines[0].strip() else "Slide"
                except Exception:
                    pass
                try:
                    tf = slide.placeholders[1].text_frame
                    body = lines[1:] if len(lines) > 1 else []
                    if body:
                        tf.text = body[0]
                        for ln in body[1:]:
                            tf.add_paragraph().text = ln
                except Exception:
                    pass
            prs.save(path)

        else:
            print(json.dumps({"ok": False, "error": f"unsupported format: {fmt}"})); return

        ok = os.path.exists(path)
        print(json.dumps({"ok": ok, "path": path, "size": os.path.getsize(path) if ok else 0}))
    except Exception as e:
        print(json.dumps({"ok": False, "error": str(e)[:200]}))

main()
